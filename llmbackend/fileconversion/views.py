from django.shortcuts import render
import os
import io
import fitz  # PyMuPDF for PDFs
from PIL import Image
import re
import string
import spacy
from nltk.corpus import wordnet
from docx import Document
from pptx import Presentation
import olefile
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt

# Load spaCy model
nlp = spacy.load("en_core_web_sm")

def has_meaning(word):
    return bool(wordnet.synsets(word))

def is_roman_numeral(text):
    roman_numerals = {'I', 'II', 'III', 'IV', 'V', 'VI', 'VII', 'VIII', 'IX', 'X',
                       'XI', 'XII', 'XIII', 'XIV', 'XV', 'XVI', 'XVII', 'XVIII', 'XIX', 'XX'}
    return text in roman_numerals

def is_heading(text, font_size, is_italic):
    doc = nlp(text)
    heading_keywords = {"introduction", "conclusion", "summary", "abstract", "chapter", "table of contents"}
    clean_text = text.translate(str.maketrans('', '', string.punctuation)).lower()
    is_heading_keyword = any(keyword in clean_text for keyword in heading_keywords)
    
    words = clean_text.split()
    meaningful_words = [word for word in words if has_meaning(word)]
    
    is_meaningful = len(meaningful_words) / len(words) > 0.5 if words else False
    
    return (len(text) < 100 and (text.istitle() or text.isupper() or font_size > 12) and not is_roman_numeral(text) and is_meaningful) or is_heading_keyword

def preprocess_text(text_blocks):
    cleaned_lines = []
    previous_heading = None
    
    for block in text_blocks:
        text = block["text"]
        font_size = block["size"]
        is_italic = block.get("italic", False)
        text = re.sub(r'\s+', ' ', text).strip()
        if not text:
            continue
        if is_heading(text, font_size, is_italic):
            if text != previous_heading:
                cleaned_lines.append(f"\n# {text}\n")
                previous_heading = text
        else:
            cleaned_lines.append(text)
            previous_heading = None
            
    return ' '.join(cleaned_lines)

def extract_text_from_pdf(pdf_file):
    try:
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        text_blocks = []
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            blocks = page.get_text("dict")["blocks"]
            for b in blocks:
                if "lines" in b:
                    for l in b["lines"]:
                        for s in l["spans"]:
                            text_blocks.append({
                                "text": s["text"],
                                "size": s["size"],
                                "italic": "italic" in s["font"]
                            })
        return preprocess_text(text_blocks)
    except Exception as e:
        return JsonResponse({"error": f"Failed to process PDF: {str(e)}"}, status=500)

def extract_text_from_docx(docx_file):
    try:
        doc = Document(docx_file)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return preprocess_text([{"text": p, "size": 12, "italic": False} for p in paragraphs])
    except Exception as e:
        return JsonResponse({"error": f"Failed to process DOCX: {str(e)}"}, status=500)

def extract_text_from_pptx(pptx_file):
    try:
        prs = Presentation(pptx_file)
        paragraphs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    paragraphs.append(shape.text)
        return preprocess_text([{"text": p, "size": 12, "italic": False} for p in paragraphs])
    except Exception as e:
        return JsonResponse({"error": f"Failed to process PPTX: {str(e)}"}, status=500)

def extract_text_from_ppt(ppt_file):
    try:
        if not olefile.isOleFile(ppt_file):
            return "Invalid OLE file"
        
        ole = olefile.OleFileIO(ppt_file)
        text_blocks = []
        for stream in ole.listdir():
            if stream[0] == 'PowerPoint Document':
                data = ole.openstream(stream)
                text = data.read().decode('utf-8', errors='ignore')
                text_blocks.append({"text": text, "size": 12, "italic": False})
        
        return preprocess_text(text_blocks)
    except Exception as e:
        return JsonResponse({"error": f"Failed to process PPT: {str(e)}"}, status=500)

@csrf_exempt
def markdown_conversion(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({"error": "No file provided"}, status=400)
        
        uploaded_file = request.FILES['file']
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        
        if file_ext == ".pdf":
            content = extract_text_from_pdf(uploaded_file)
        elif file_ext == ".docx":
            content = extract_text_from_docx(uploaded_file)
        elif file_ext == ".pptx":
            content = extract_text_from_pptx(uploaded_file)
        elif file_ext == ".ppt":
            content = extract_text_from_ppt(uploaded_file)
        else:
            return JsonResponse({"error": "Unsupported file format"}, status=400)

        if isinstance(content, JsonResponse):
            return content
        
        return HttpResponse(content, content_type='text/markdown')
    else:
        return JsonResponse({"error": "Invalid request method"}, status=405)
